# -*- coding: utf-8 -*-
"""
生成中期答辩PPT
基于文字脚本自动生成PowerPoint演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """创建PPT演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 定义配色方案
    colors = {
        'primary': RGBColor(0, 102, 255),      # 深蓝 #0066FF
        'secondary': RGBColor(0, 196, 140),    # 青绿 #00C48C
        'background': RGBColor(30, 34, 45),    # 深灰 #1E222D
        'text': RGBColor(209, 212, 220),       # 浅灰 #D1D4DC
        'white': RGBColor(255, 255, 255)
    }
    
    # 幻灯片1: 封面页
    slide1 = add_title_slide(prs, colors)
    
    # 幻灯片2: 项目背景与研究意义
    slide2 = add_background_slide(prs, colors)
    
    # 幻灯片3: 技术路线与创新点
    slide3 = add_technical_architecture_slide(prs, colors)
    
    # 幻灯片4: 核心工作内容
    slide4 = add_core_work_slide(prs, colors)
    
    # 幻灯片5: 实施进展与时间节点
    slide5 = add_progress_slide(prs, colors)
    
    # 幻灯片6: 核心成果展示（数据篇）
    slide6 = add_data_results_slide(prs, colors)
    
    # 幻灯片7: 核心成果展示（系统篇）
    slide7 = add_system_results_slide(prs, colors)
    
    # 幻灯片8: 系统演示
    slide8 = add_demo_slide(prs, colors)
    
    # 幻灯片9: 团队分工与协作
    slide9 = add_team_slide(prs, colors)
    
    # 幻灯片10: 遇到的困难与解决方案
    slide10 = add_challenges_slide(prs, colors)
    
    # 幻灯片11: 下一步工作计划
    slide11 = add_future_work_slide(prs, colors)
    
    # 幻灯片12: 总结与致谢
    slide12 = add_summary_slide(prs, colors)
    
    # 幻灯片13: Q&A准备
    slide13 = add_qa_slide(prs, colors)
    
    return prs


def add_title_slide(prs, colors):
    """幻灯片1: 封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "基于因果关系的金融资产收益率预测系统"
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "中期答辩报告"
    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = colors['secondary']
    
    # 项目信息
    info_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(2))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True
    
    info_text = [
        "项目类型: 创新训练项目",
        "负责人: 刘一鸣 (23281172)",
        "成员: 彭闻聪、吴伦凯",
        "指导教师: 刘海洋",
        "答辩日期: 2025年10月"
    ]
    
    for i, text in enumerate(info_text):
        if i > 0:
            info_frame.add_paragraph()
        p = info_frame.paragraphs[i]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(8)
    
    return slide


def add_background_slide(prs, colors):
    """幻灯片2: 项目背景与研究意义"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "为什么需要因果关系?"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 左侧: 问题痛点
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = "问题痛点"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = colors['secondary']
    p.space_after = Pt(12)
    
    problems = [
        "• 传统模型依赖统计相关性",
        "• 市场突变时预测失效",
        "• 黑箱模型缺乏可解释性",
        "• 无法捕捉股票间联动效应"
    ]
    
    for prob in problems:
        left_frame.add_paragraph()
        p = left_frame.paragraphs[-1]
        p.text = prob
        p.font.size = Pt(18)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(10)
    
    # 右侧: 解决方案
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4.5), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = "解决方案"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = colors['secondary']
    p.space_after = Pt(12)
    
    solutions = [
        "• 引入因果推理机制",
        "• 构建股票影响关系网络",
        "• 提供可解释的预测路径",
        "• 多模态数据融合"
    ]
    
    for sol in solutions:
        right_frame.add_paragraph()
        p = right_frame.paragraphs[-1]
        p.text = sol
        p.font.size = Pt(18)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(10)
    
    return slide


def add_technical_architecture_slide(prs, colors):
    """幻灯片3: 技术路线与创新点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "ICSFP平台技术架构"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 流程图
    flow_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
    flow_frame = flow_box.text_frame
    flow_frame.word_wrap = True
    
    p = flow_frame.paragraphs[0]
    p.text = "数据输入 → 因果发现 → 模型预测 → 结果输出"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)
    
    flow_frame.add_paragraph()
    p = flow_frame.paragraphs[1]
    p.text = "价格/文本    Granger/CUTS+    StockNet/MSGNet    可视化/API"
    p.font.size = Pt(16)
    p.font.color.rgb = colors['text']
    p.alignment = PP_ALIGN.CENTER
    
    # 三大创新点
    innovations_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(4))
    innovations_frame = innovations_box.text_frame
    innovations_frame.word_wrap = True
    
    p = innovations_frame.paragraphs[0]
    p.text = "三大创新点"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = colors['secondary']
    p.space_after = Pt(12)
    
    innovations = [
        "1. 因果嵌入深度学习",
        "   • Granger因果检验 + CUTS+算法",
        "   • 生成31×31因果邻接矩阵",
        "   • 图神经网络编码股票关系",
        "",
        "2. 多模态数据融合",
        "   • 价格序列(VMD分解) + 社交媒体文本(GloVe嵌入) + 因果图特征",
        "",
        "3. 系统工程化实现",
        "   • RESTful API(8个端点) + Web可视化界面 + Docker容器化部署"
    ]
    
    for innov in innovations:
        innovations_frame.add_paragraph()
        p = innovations_frame.paragraphs[-1]
        p.text = innov
        p.font.size = Pt(16) if innov.startswith("   ") else Pt(18)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(6)
    
    return slide


def add_core_work_slide(prs, colors):
    """幻灯片4: 核心工作内容"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "项目主要工作与成果"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 工作内容
    work_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    work_frame = work_box.text_frame
    work_frame.word_wrap = True
    
    work_items = [
        ("🔹 系统工程化实现", [
            "RESTful API服务 (2000行)",
            "现代化Web前端界面 (1400行)",
            "完整的Docker部署方案"
        ]),
        ("🔹 性能优化与重构", [
            "GPU加速因果图计算 (800行)",
            "模块化架构设计 (600行)",
            "代码结构优化与重组"
        ]),
        ("🔹 评估与测试体系", [
            "完整评估指标系统 (400行)",
            "数据预处理流程 (500行)",
            "实验验证与性能分析"
        ]),
        ("🔹 项目文档与管理", [
            "技术文档15篇 (API文档、部署指南等)",
            "GitHub版本管理",
            "完整的开发规范"
        ])
    ]
    
    first = True
    for title, items in work_items:
        if not first:
            work_frame.add_paragraph()
        first = False
        
        p = work_frame.paragraphs[-1]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = colors['secondary']
        p.space_after = Pt(8)
        
        for item in items:
            work_frame.add_paragraph()
            p = work_frame.paragraphs[-1]
            p.text = f"  • {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = colors['text']
            p.space_after = Pt(4)
    
    # 代码统计
    work_frame.add_paragraph()
    work_frame.add_paragraph()
    p = work_frame.paragraphs[-1]
    p.text = "代码统计: 总计约8,000行代码"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_progress_slide(prs, colors):
    """幻灯片5: 实施进展与时间节点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "项目进度一览"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 时间轴
    timeline_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    timeline_frame = timeline_box.text_frame
    timeline_frame.word_wrap = True
    
    phases = [
        ("✅ 2025.06-2025.08: 文献调研与数据采集", "完成度: 100% | 成果: 88只股票数据集"),
        ("✅ 2025.08-2025.11: 系统架构与算法实现", "完成度: 100% | 成果: API服务、前端、因果模块"),
        ("🔄 2025.11-2026.03: 模型优化与接口开发", "完成度: 30% (进行中) | 任务: 超参调优、实时数据"),
        ("📅 2026.03-2026.05: 系统测试与性能验证", "完成度: 0% (计划中) | 任务: 回测、论文撰写")
    ]
    
    first = True
    for phase, detail in phases:
        if not first:
            timeline_frame.add_paragraph()
            timeline_frame.add_paragraph()
        first = False
        
        p = timeline_frame.paragraphs[-1]
        p.text = phase
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = colors['white']
        p.space_after = Pt(8)
        
        timeline_frame.add_paragraph()
        p = timeline_frame.paragraphs[-1]
        p.text = f"  {detail}"
        p.font.size = Pt(16)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(4)
    
    return slide


def add_data_results_slide(prs, colors):
    """幻灯片6: 核心成果展示(数据篇)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "量化指标与实验结果"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 左侧: 模型性能表格
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(3))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = "模型性能"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = colors['secondary']
    p.space_after = Pt(10)
    
    table_data = [
        "数据集          准确率    MCC",
        "━━━━━━━━━━━━━━━━━━━━━",
        "CMIN-CN      61.5%   0.230",
        "ACL18        58.2%   0.165",
        "随机基线      50.0%   0.000"
    ]
    
    for row in table_data:
        left_frame.add_paragraph()
        p = left_frame.paragraphs[-1]
        p.text = row
        p.font.name = 'Consolas'
        p.font.size = Pt(14) if "━" in row else Pt(16)
        p.font.color.rgb = colors['white'] if "CMIN-CN" in row else colors['text']
        p.font.bold = "CMIN-CN" in row
        p.space_after = Pt(4)
    
    left_frame.add_paragraph()
    left_frame.add_paragraph()
    p = left_frame.paragraphs[-1]
    p.text = "性能分析:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = colors['secondary']
    p.space_after = Pt(6)
    
    analysis = [
        "• 提升11.5个百分点",
        "• MCC=0.230有实际价值",
        "• 60%+属良好水平"
    ]
    
    for item in analysis:
        left_frame.add_paragraph()
        p = left_frame.paragraphs[-1]
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(4)
    
    # 右侧: 系统功能
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4.5), Inches(3))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = "系统功能"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = colors['secondary']
    p.space_after = Pt(10)
    
    features = [
        "• API端点: 8个",
        "• 支持股票: 88只",
        "• 词表规模: 29,867个",
        "• 因果图维度: 31×31",
        "• 代码规模: 8,000行"
    ]
    
    for feature in features:
        right_frame.add_paragraph()
        p = right_frame.paragraphs[-1]
        p.text = feature
        p.font.size = Pt(18)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(10)
    
    return slide


def add_system_results_slide(prs, colors):
    """幻灯片7: 核心成果展示(系统篇)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "软件系统与工程实现"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 四宫格布局
    quadrants = [
        (Inches(0.5), Inches(1.2), "Web界面", [
            "• 专业暗色主题",
            "• Chart.js动态图表",
            "• 交互式预测表单",
            "• 实时结果展示"
        ]),
        (Inches(5.5), Inches(1.2), "API服务", [
            "• Flask 3.1.2框架",
            "• RESTful规范",
            "• 8个核心端点",
            "• 完整错误处理"
        ]),
        (Inches(0.5), Inches(4.2), "因果图可视化", [
            "• 31×31邻接矩阵",
            "• GPU加速计算",
            "• 影响力分析",
            "• 路径追踪"
        ]),
        (Inches(5.5), Inches(4.2), "部署方案", [
            "• Docker容器化",
            "• docker-compose编排",
            "• 一键启动脚本",
            "• 完整技术文档"
        ])
    ]
    
    for left, top, title, items in quadrants:
        box = slide.shapes.add_textbox(left, top, Inches(4.5), Inches(2.8))
        frame = box.text_frame
        frame.word_wrap = True
        
        p = frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = colors['secondary']
        p.space_after = Pt(8)
        
        for item in items:
            frame.add_paragraph()
            p = frame.paragraphs[-1]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = colors['text']
            p.space_after = Pt(6)
    
    return slide


def add_demo_slide(prs, colors):
    """幻灯片8: 系统演示"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "系统实时演示"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 演示流程
    demo_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    demo_frame = demo_box.text_frame
    demo_frame.word_wrap = True
    
    steps = [
        ("步骤1: 首页访问", [
            "URL: http://127.0.0.1:5000",
            "显示统计面板、预测表单、图表区域"
        ]),
        ("步骤2: 单股票预测", [
            "选择股票: AAPL",
            "日期范围: 2015-10-01 至 2015-10-05",
            "勾选: 使用因果图 → 点击预测"
        ]),
        ("步骤3: 结果展示", [
            "折线图: UP/DOWN概率曲线",
            "数据表: 每日预测结果",
            "置信度条: 可视化概率分布"
        ]),
        ("步骤4: 因果图查看", [
            "访问关于页面",
            "查看因果影响路径和技术架构说明"
        ])
    ]
    
    first = True
    for step, details in steps:
        if not first:
            demo_frame.add_paragraph()
        first = False
        
        p = demo_frame.paragraphs[-1]
        p.text = step
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = colors['secondary']
        p.space_after = Pt(8)
        
        for detail in details:
            demo_frame.add_paragraph()
            p = demo_frame.paragraphs[-1]
            p.text = f"  • {detail}"
            p.font.size = Pt(16)
            p.font.color.rgb = colors['text']
            p.space_after = Pt(4)
    
    # 提示: 准备截图
    demo_frame.add_paragraph()
    demo_frame.add_paragraph()
    p = demo_frame.paragraphs[-1]
    p.text = "📸 提示: 请在实际演示时插入系统截图"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = colors['secondary']
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_team_slide(prs, colors):
    """幻灯片9: 团队分工与协作"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "团队成员贡献"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 三栏布局
    members = [
        (Inches(0.5), "刘一鸣 (负责人)", "系统架构与API开发", [
            "Flask REST API (2000行)",
            "Web前端界面 (1400行)",
            "项目管理与文档"
        ], "每周15小时, 共180小时"),
        (Inches(3.7), "彭闻聪", "因果算法与优化", [
            "GPU加速因果计算 (800行)",
            "模块化架构重构 (600行)",
            "因果分析接口"
        ], "每周12小时, 共144小时"),
        (Inches(6.9), "吴伦凯", "数据处理与模型训练", [
            "数据预处理流程 (500行)",
            "评估指标体系 (400行)",
            "实验与性能测试"
        ], "每周10小时, 共120小时")
    ]
    
    for left, name, role, achievements, time in members:
        box = slide.shapes.add_textbox(left, Inches(1.2), Inches(3), Inches(5.8))
        frame = box.text_frame
        frame.word_wrap = True
        
        p = frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors['secondary']
        p.space_after = Pt(8)
        
        frame.add_paragraph()
        p = frame.paragraphs[-1]
        p.text = f"职责: {role}"
        p.font.size = Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(8)
        
        frame.add_paragraph()
        p = frame.paragraphs[-1]
        p.text = "成果:"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = colors['white']
        p.space_after = Pt(6)
        
        for achievement in achievements:
            frame.add_paragraph()
            p = frame.paragraphs[-1]
            p.text = f"• {achievement}"
            p.font.size = Pt(12)
            p.font.color.rgb = colors['text']
            p.space_after = Pt(4)
        
        frame.add_paragraph()
        frame.add_paragraph()
        p = frame.paragraphs[-1]
        p.text = f"投入: {time}"
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = colors['text']
    
    return slide


def add_challenges_slide(prs, colors):
    """幻灯片10: 遇到的困难与解决方案"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "技术难点与解决思路"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 四个挑战
    challenges_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    challenges_frame = challenges_box.text_frame
    challenges_frame.word_wrap = True
    
    challenges = [
        ("1. 因果图计算复杂性", 
         "问题: 88×88因果图计算耗时长",
         "解决: GPU并行 + 缓存机制",
         "效果: 时间从小时级降到分钟级"),
        ("2. 多源数据时序对齐",
         "问题: 文本与价格数据时间戳不一致",
         "解决: 统一时间映射 + 滑动窗口",
         "效果: 成功对齐88只股票数据"),
        ("3. 模型可解释性",
         "问题: 深度神经网络黑箱",
         "解决: 注意力机制 + 因果归因",
         "效果: 提供可视化因果路径"),
        ("4. 金融预测难度",
         "问题: 市场受多因素影响",
         "认识: 60%准确率已属良好",
         "重点: 可解释性比准确率更重要")
    ]
    
    first = True
    for title, problem, solution, result in challenges:
        if not first:
            challenges_frame.add_paragraph()
        first = False
        
        p = challenges_frame.paragraphs[-1]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors['secondary']
        p.space_after = Pt(6)
        
        for text in [problem, solution, result]:
            challenges_frame.add_paragraph()
            p = challenges_frame.paragraphs[-1]
            p.text = f"  {text}"
            p.font.size = Pt(14)
            p.font.color.rgb = colors['text']
            p.space_after = Pt(4)
    
    return slide


def add_future_work_slide(prs, colors):
    """幻灯片11: 下一步工作计划"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "后续研究方向"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 工作计划
    plan_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    plan_frame = plan_box.text_frame
    plan_frame.word_wrap = True
    
    plans = [
        ("第三阶段 (2025.11-2026.03): 优化与扩展", [
            "超参数网格搜索",
            "扩展数据集规模",
            "接入实时数据流",
            "API性能压力测试"
        ]),
        ("第四阶段 (2026.03-2026.05): 验证与部署", [
            "跨市场验证 (A股、港股)",
            "回测系统开发",
            "云平台部署",
            "用户体验测试"
        ]),
        ("学术产出计划", [
            "整理消融实验与对比研究",
            "撰写会议论文 (目标: CCF-C类)",
            "完成项目结题报告",
            "开源代码与文档"
        ])
    ]
    
    first = True
    for title, items in plans:
        if not first:
            plan_frame.add_paragraph()
        first = False
        
        p = plan_frame.paragraphs[-1]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = colors['secondary']
        p.space_after = Pt(8)
        
        for item in items:
            plan_frame.add_paragraph()
            p = plan_frame.paragraphs[-1]
            p.text = f"  • {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = colors['text']
            p.space_after = Pt(6)
    
    return slide


def add_summary_slide(prs, colors):
    """幻灯片12: 总结与致谢"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "中期总结"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # 左侧: 已完成成果
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(4))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = "已完成成果"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = colors['secondary']
    p.space_after = Pt(10)
    
    achievements = [
        "✅ 完整技术栈实现",
        "✅ 8个API端点",
        "✅ 现代化Web界面",
        "✅ 因果图GPU加速",
        "✅ 88只股票数据集",
        "✅ ACC 61.5%, MCC 0.230",
        "✅ 8,000行代码 + 15篇文档"
    ]
    
    for item in achievements:
        left_frame.add_paragraph()
        p = left_frame.paragraphs[-1]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(8)
    
    # 右侧: 核心贡献
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4.5), Inches(4))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = "核心贡献"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = colors['secondary']
    p.space_after = Pt(10)
    
    contributions = [
        "• 因果推理 + 深度学习融合",
        "• 系统工程化实现",
        "• 可解释性分析",
        "• 完整的软件交付"
    ]
    
    for item in contributions:
        right_frame.add_paragraph()
        p = right_frame.paragraphs[-1]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(12)
    
    # 致谢
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.word_wrap = True
    
    thanks_text = [
        "感谢指导老师刘海洋的悉心指导",
        "感谢计算机学院提供的实验环境",
        "感谢HCSF开源项目提供的基础框架"
    ]
    
    for text in thanks_text:
        if thanks_frame.paragraphs[0].text:
            thanks_frame.add_paragraph()
        p = thanks_frame.paragraphs[-1]
        p.text = text
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = colors['text']
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)
    
    # 结束语
    ending_box = slide.shapes.add_textbox(Inches(2), Inches(6.8), Inches(6), Inches(0.5))
    ending_frame = ending_box.text_frame
    p = ending_frame.paragraphs[0]
    p.text = "以上是我们的中期答辩汇报, 请各位老师批评指正!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_qa_slide(prs, colors):
    """幻灯片13: Q&A准备"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Q&A 准备"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # Q&A内容
    qa_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    qa_frame = qa_box.text_frame
    qa_frame.word_wrap = True
    
    qas = [
        ("Q1: 为什么准确率只有61.5%?",
         "A: 金融预测本质难度高, 60%+已属良好水平, 重点在可解释性"),
        ("Q2: 因果图如何嵌入模型?",
         "A: 通过图神经网络将邻接矩阵编码为特征, 与价格/文本特征融合"),
        ("Q3: 与StockNet等基线比较?",
         "A: 我们复现了StockNet并加入因果增强, CMIN-CN上有提升"),
        ("Q4: 系统实用价值?",
         "A: 提供可解释预测、验证因果+深度学习可行性、教育价值"),
        ("Q5: 你们的主要贡献?",
         "A: 把学术算法工程化成可部署平台 (API+前端+因果优化)")
    ]
    
    first = True
    for question, answer in qas:
        if not first:
            qa_frame.add_paragraph()
        first = False
        
        p = qa_frame.paragraphs[-1]
        p.text = question
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors['secondary']
        p.space_after = Pt(4)
        
        qa_frame.add_paragraph()
        p = qa_frame.paragraphs[-1]
        p.text = answer
        p.font.size = Pt(14)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(6)
    
    return slide


def main():
    """主函数"""
    print("开始生成PPT...")
    
    # 创建演示文稿
    prs = create_presentation()
    
    # 保存文件
    output_file = "中期答辩_ICSFP.pptx"
    prs.save(output_file)
    
    print(f"✅ PPT生成成功: {output_file}")
    print(f"   共 {len(prs.slides)} 张幻灯片")
    print("\n提示:")
    print("  1. 请在幻灯片8中插入系统截图")
    print("  2. 可根据需要调整字体大小和布局")
    print("  3. 建议使用PowerPoint打开后进一步美化")
    print("  4. 记得导出PDF备份版本")


if __name__ == "__main__":
    main()
